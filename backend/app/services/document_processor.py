"""
Document processing service using LangChain + EasyOCR
"""
import os
from typing import Optional, List, Dict
from pathlib import Path
import logging
import csv as csv_module
import email
from email.parser import BytesParser
from email.policy import default
import re
import zipfile
import xml.etree.ElementTree as ET

from langchain.document_loaders import PyPDFLoader, Docx2txtLoader, TextLoader
from pdf2image import convert_from_path
import easyocr
from PIL import Image
import numpy as np
# Lazy imports for optional dependencies - import only when needed

from app.core.config import settings

logger = logging.getLogger(__name__)


class DocumentProcessor:
    """Process documents using LangChain loaders with EasyOCR for OCR"""
    
    def __init__(self):
        # Initialize EasyOCR reader (lazy load on first use)
        self._easyocr_reader = None
    
    @property
    def easyocr_reader(self):
        """Lazy load EasyOCR reader"""
        if self._easyocr_reader is None:
            try:
                logger.info("Initializing EasyOCR reader...")
                self._easyocr_reader = easyocr.Reader(['en'], gpu=False)  # Set gpu=True if you have GPU
                logger.info("EasyOCR reader initialized")
            except Exception as e:
                logger.error(f"Failed to initialize EasyOCR: {e}", exc_info=True)
                # Create a dummy reader that will fail gracefully
                raise RuntimeError(f"EasyOCR initialization failed: {str(e)}")
        return self._easyocr_reader
    
    def process_document(self, file_path: str, file_type: str) -> Dict:
        """
        Process a document and extract text with metadata
        
        Returns:
            dict with keys: text, page_count, requires_ocr, pages (list of page texts)
        """
        try:
            file_ext = file_type.lower().lstrip('.')
            
            if file_ext == "pdf":
                return self._process_pdf(file_path)
            elif file_ext in ["docx", "doc"]:
                return self._process_docx(file_path)
            elif file_ext == "txt":
                return self._process_txt(file_path)
            elif file_ext in ["jpg", "jpeg", "png", "gif", "bmp", "webp", "tiff", "tif"]:
                return self._process_image(file_path)
            elif file_ext == "xlsx":
                return self._process_xlsx(file_path)
            elif file_ext == "csv":
                return self._process_csv(file_path)
            elif file_ext == "msg":
                return self._process_msg(file_path)
            elif file_ext == "eml":
                return self._process_eml(file_path)
            elif file_ext == "pptx":
                return self._process_pptx(file_path)
            elif file_ext == "odt":
                return self._process_odt(file_path)
            elif file_ext == "ods":
                return self._process_ods(file_path)
            elif file_ext == "epub":
                return self._process_epub(file_path)
            elif file_ext == "xps":
                return self._process_xps(file_path)
            else:
                raise ValueError(f"Unsupported file type: {file_ext}")
        except Exception as e:
            logger.error(f"Error processing document {file_path}: {str(e)}", exc_info=True)
            raise
    
    def _process_pdf(self, file_path: str) -> Dict:
        """Process PDF file using LangChain with EasyOCR fallback"""
        try:
            # Try LangChain PyPDFLoader first
            loader = PyPDFLoader(file_path)
            documents = loader.load()
            
            pages = []
            full_text = ""
            requires_ocr = False
            
            for i, doc in enumerate(documents, 1):
                page_text = doc.page_content.strip()
                
                # Check if page has meaningful text
                if len(page_text) < 50:  # Likely scanned or empty
                    requires_ocr = True
                    logger.info(f"Page {i} appears to be scanned, performing OCR...")
                    
                    try:
                        # Convert PDF page to image and OCR
                        images = convert_from_path(file_path, first_page=i, last_page=i, dpi=300)
                        if images:
                            ocr_text = self._ocr_image(images[0])
                            page_text = ocr_text
                            method = "ocr"
                        else:
                            method = "ocr_needed"
                    except Exception as e:
                        logger.warning(f"OCR failed for page {i}: {e}")
                        method = "ocr_needed"
                else:
                    method = "extraction"
                
                pages.append({
                    "page_number": i,
                    "text": page_text,
                    "method": method
                })
                full_text += page_text + "\n\n"
            
            return {
                "text": full_text.strip(),
                "page_count": len(pages),
                "requires_ocr": requires_ocr,
                "pages": pages
            }
            
        except Exception as e:
            logger.error(f"Error processing PDF {file_path}: {e}", exc_info=True)
            raise
    
    def _process_docx(self, file_path: str) -> Dict:
        """Process DOCX file using LangChain"""
        try:
            loader = Docx2txtLoader(file_path)
            documents = loader.load()
            
            full_text = "\n\n".join([doc.page_content for doc in documents])
            word_count = len(full_text.split())
            estimated_pages = max(1, word_count // 500)
            
            return {
                "text": full_text,
                "page_count": estimated_pages,
                "requires_ocr": False,
                "pages": [{"page_number": 1, "text": full_text, "method": "extraction"}]
            }
        except Exception as e:
            logger.error(f"Error processing DOCX {file_path}: {e}", exc_info=True)
            raise
    
    def _process_txt(self, file_path: str) -> Dict:
        """Process TXT file using LangChain"""
        try:
            loader = TextLoader(file_path, encoding='utf-8')
            documents = loader.load()
            
            full_text = "\n\n".join([doc.page_content for doc in documents])
            word_count = len(full_text.split())
            estimated_pages = max(1, word_count // 500)
            
            return {
                "text": full_text,
                "page_count": estimated_pages,
                "requires_ocr": False,
                "pages": [{"page_number": 1, "text": full_text, "method": "extraction"}]
            }
        except Exception as e:
            logger.error(f"Error processing TXT {file_path}: {e}", exc_info=True)
            raise
    
    def _process_image(self, file_path: str) -> Dict:
        """Process image file using EasyOCR"""
        try:
            logger.info(f"Performing OCR on image: {file_path}")
            img = Image.open(file_path)
            
            # Perform OCR using EasyOCR
            ocr_text = self._ocr_image(img)
            
            return {
                "text": ocr_text,
                "page_count": 1,
                "requires_ocr": True,
                "pages": [{"page_number": 1, "text": ocr_text, "method": "ocr"}]
            }
        except Exception as e:
            logger.error(f"Error processing image {file_path}: {e}", exc_info=True)
            raise
    
    def _ocr_image(self, image) -> str:
        """Perform OCR on an image using EasyOCR"""
        try:
            # Convert PIL Image to numpy array if needed
            if hasattr(image, 'save'):  # PIL Image
                img_array = np.array(image)
            else:
                img_array = image
            
            # Perform OCR
            results = self.easyocr_reader.readtext(img_array)
            
            # Combine all detected text
            ocr_text = "\n".join([result[1] for result in results])
            
            return ocr_text
        except Exception as e:
            logger.error(f"EasyOCR error: {e}", exc_info=True)
            return ""
    
    def _process_xlsx(self, file_path: str) -> Dict:
        """Process Excel file (.xlsx)"""
        try:
            from openpyxl import load_workbook
            workbook = load_workbook(file_path, data_only=True)
            full_text = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                full_text.append(f"\n\n=== Sheet: {sheet_name} ===\n\n")
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        full_text.append(row_text)
            
            full_text_str = "\n".join(full_text)
            word_count = len(full_text_str.split())
            estimated_pages = max(1, word_count // 500)
            
            return {
                "text": full_text_str,
                "page_count": estimated_pages,
                "requires_ocr": False,
                "pages": [{"page_number": 1, "text": full_text_str, "method": "extraction"}]
            }
        except Exception as e:
            logger.error(f"Error processing XLSX {file_path}: {e}", exc_info=True)
            raise

    def _process_csv(self, file_path: str) -> Dict:
        """Process CSV file"""
        try:
            full_text = []
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                reader = csv_module.reader(f)
                for row in reader:
                    row_text = "\t".join(row)
                    if row_text.strip():
                        full_text.append(row_text)
            
            full_text_str = "\n".join(full_text)
            word_count = len(full_text_str.split())
            estimated_pages = max(1, word_count // 500)
            
            return {
                "text": full_text_str,
                "page_count": estimated_pages,
                "requires_ocr": False,
                "pages": [{"page_number": 1, "text": full_text_str, "method": "extraction"}]
            }
        except Exception as e:
            logger.error(f"Error processing CSV {file_path}: {e}", exc_info=True)
            raise

    def _process_msg(self, file_path: str) -> Dict:
        """Process Outlook .msg file"""
        try:
            import extract_msg
            msg = extract_msg.Message(file_path)
            
            # Extract email fields
            text_parts = []
            text_parts.append(f"From: {msg.sender}")
            text_parts.append(f"To: {msg.to}")
            text_parts.append(f"CC: {msg.cc if hasattr(msg, 'cc') else ''}")
            text_parts.append(f"Subject: {msg.subject}")
            text_parts.append(f"Date: {msg.date}")
            text_parts.append("\n--- Body ---\n")
            text_parts.append(msg.body)
            
            # Attachments info
            if msg.attachments:
                text_parts.append("\n--- Attachments ---\n")
                for att in msg.attachments:
                    text_parts.append(f"{att.shortFilename}")
            
            full_text = "\n".join(text_parts)
            word_count = len(full_text.split())
            estimated_pages = max(1, word_count // 500)
            
            return {
                "text": full_text,
                "page_count": estimated_pages,
                "requires_ocr": False,
                "pages": [{"page_number": 1, "text": full_text, "method": "extraction"}]
            }
        except Exception as e:
            logger.error(f"Error processing MSG {file_path}: {e}", exc_info=True)
            raise

    def _process_eml(self, file_path: str) -> Dict:
        """Process .eml email file"""
        try:
            with open(file_path, 'rb') as f:
                msg = BytesParser(policy=default).parse(f)
            
            text_parts = []
            text_parts.append(f"From: {msg['From']}")
            text_parts.append(f"To: {msg['To']}")
            text_parts.append(f"CC: {msg.get('CC', '')}")
            text_parts.append(f"Subject: {msg['Subject']}")
            text_parts.append(f"Date: {msg['Date']}")
            text_parts.append("\n--- Body ---\n")
            
            # Get body text
            body = ""
            if msg.is_multipart():
                for part in msg.walk():
                    content_type = part.get_content_type()
                    if content_type == "text/plain":
                        body = part.get_payload(decode=True).decode('utf-8', errors='ignore')
                        break
                    elif content_type == "text/html":
                        # Fallback to HTML if plain text not available
                        if not body:
                            body = part.get_payload(decode=True).decode('utf-8', errors='ignore')
            else:
                body = msg.get_payload(decode=True).decode('utf-8', errors='ignore')
            
            text_parts.append(body)
            full_text = "\n".join(text_parts)
            word_count = len(full_text.split())
            estimated_pages = max(1, word_count // 500)
            
            return {
                "text": full_text,
                "page_count": estimated_pages,
                "requires_ocr": False,
                "pages": [{"page_number": 1, "text": full_text, "method": "extraction"}]
            }
        except Exception as e:
            logger.error(f"Error processing EML {file_path}: {e}", exc_info=True)
            raise

    def _process_pptx(self, file_path: str) -> Dict:
        """Process PowerPoint .pptx file"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            full_text = []
            pages = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                slide_text.append(f"--- Slide {slide_num} ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                slide_text_str = "\n".join(slide_text)
                full_text.append(slide_text_str)
                pages.append({
                    "page_number": slide_num,
                    "text": slide_text_str,
                    "method": "extraction"
                })
            
            full_text_str = "\n\n".join(full_text)
            page_count = len(prs.slides)
            
            return {
                "text": full_text_str,
                "page_count": page_count,
                "requires_ocr": False,
                "pages": pages
            }
        except Exception as e:
            logger.error(f"Error processing PPTX {file_path}: {e}", exc_info=True)
            raise

    def _process_odt(self, file_path: str) -> Dict:
        """Process OpenDocument Text (.odt)"""
        try:
            from odf.opendocument import load as odf_load
            from odf.text import P as OdfP
            doc = odf_load(file_path)
            full_text = []
            
            for paragraph in doc.getElementsByType(OdfP):
                text = paragraph.getTextContent()
                if text.strip():
                    full_text.append(text)
            
            full_text_str = "\n\n".join(full_text)
            word_count = len(full_text_str.split())
            estimated_pages = max(1, word_count // 500)
            
            return {
                "text": full_text_str,
                "page_count": estimated_pages,
                "requires_ocr": False,
                "pages": [{"page_number": 1, "text": full_text_str, "method": "extraction"}]
            }
        except Exception as e:
            logger.error(f"Error processing ODT {file_path}: {e}", exc_info=True)
            raise

    def _process_ods(self, file_path: str) -> Dict:
        """Process OpenDocument Spreadsheet (.ods)"""
        try:
            from odf.opendocument import load as odf_load
            from odf.table import Table, TableRow, TableCell
            doc = odf_load(file_path)
            full_text = []
            
            for table in doc.getElementsByType(Table):
                for row in table.getElementsByType(TableRow):
                    row_text = []
                    for cell in row.getElementsByType(TableCell):
                        cell_text = cell.getTextContent().strip()
                        row_text.append(cell_text)
                    if any(row_text):
                        full_text.append("\t".join(row_text))
            
            full_text_str = "\n".join(full_text)
            word_count = len(full_text_str.split())
            estimated_pages = max(1, word_count // 500)
            
            return {
                "text": full_text_str,
                "page_count": estimated_pages,
                "requires_ocr": False,
                "pages": [{"page_number": 1, "text": full_text_str, "method": "extraction"}]
            }
        except Exception as e:
            logger.error(f"Error processing ODS {file_path}: {e}", exc_info=True)
            raise

    def _process_epub(self, file_path: str) -> Dict:
        """Process EPUB ebook"""
        try:
            import ebooklib
            from ebooklib import epub
            book = epub.read_epub(file_path)
            full_text = []
            pages = []
            page_num = 1
            
            for item in book.get_items():
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    # Extract text from HTML content
                    content = item.get_content().decode('utf-8', errors='ignore')
                    # Simple HTML tag removal
                    text = re.sub(r'<[^>]+>', '', content)
                    text = re.sub(r'\s+', ' ', text).strip()
                    if text:
                        full_text.append(text)
                        pages.append({
                            "page_number": page_num,
                            "text": text,
                            "method": "extraction"
                        })
                        page_num += 1
            
            full_text_str = "\n\n".join(full_text)
            
            return {
                "text": full_text_str,
                "page_count": max(1, len(pages)),
                "requires_ocr": False,
                "pages": pages if pages else [{"page_number": 1, "text": full_text_str, "method": "extraction"}]
            }
        except Exception as e:
            logger.error(f"Error processing EPUB {file_path}: {e}", exc_info=True)
            raise

    def _process_xps(self, file_path: str) -> Dict:
        """Process XPS document (XML Paper Specification)"""
        try:
            # XPS is essentially a ZIP file with XML
            # For now, we'll try to extract text from XML, fallback to OCR
            text_parts = []
            try:
                with zipfile.ZipFile(file_path, 'r') as xps_zip:
                    # Try to extract text from FixedDocumentSequence
                    for name in xps_zip.namelist():
                        if name.endswith('.fpage') or 'FixedDocument' in name:
                            try:
                                content = xps_zip.read(name)
                                root = ET.fromstring(content)
                                # Extract text from XML (simplified)
                                for elem in root.iter():
                                    if elem.text and elem.text.strip():
                                        text_parts.append(elem.text.strip())
                            except Exception:
                                continue
            except Exception as e:
                logger.warning(f"XPS ZIP parsing failed: {e}")
            
            if not text_parts:
                # Fallback: treat as image and use OCR
                logger.info("XPS file has no extractable text, using OCR")
                return self._process_image(file_path)
            
            full_text = "\n".join(text_parts)
            word_count = len(full_text.split())
            estimated_pages = max(1, word_count // 500)
            
            return {
                "text": full_text,
                "page_count": estimated_pages,
                "requires_ocr": False,
                "pages": [{"page_number": 1, "text": full_text, "method": "extraction"}]
            }
        except Exception as e:
            logger.error(f"Error processing XPS {file_path}: {e}", exc_info=True)
            # Fallback to OCR if XPS parsing fails
            logger.info("XPS parsing failed, attempting OCR fallback")
            try:
                return self._process_image(file_path)
            except Exception as ocr_error:
                logger.error(f"OCR fallback also failed: {ocr_error}", exc_info=True)
                raise
    
    def chunk_text(
        self, 
        text: str, 
        chunk_size: int = None, 
        chunk_overlap: int = None,
        preserve_paragraphs: bool = True,
        page_mapping: Optional[List[Dict]] = None
    ) -> List[Dict]:
        """
        Chunk text into smaller pieces for vector search
        
        Args:
            text: Full text to chunk
            chunk_size: Maximum characters per chunk
            chunk_overlap: Overlap between chunks
            preserve_paragraphs: Try to keep paragraphs intact
            page_mapping: List of dicts with page boundaries: [{"page": 1, "start": 0, "end": 500}, ...]
        
        Returns:
            List of chunk dicts with: content, start_char, end_char, chunk_index, page_number
        """
        chunk_size = chunk_size or settings.CHUNK_SIZE
        chunk_overlap = chunk_overlap or settings.CHUNK_OVERLAP
        
        chunks = []
        
        # Helper function to find page number for a character position
        def find_page_number(char_pos: int) -> Optional[int]:
            """Find the page number for a given character position"""
            if not page_mapping:
                return None
            for page_info in page_mapping:
                if page_info["start"] <= char_pos < page_info["end"]:
                    return page_info["page"]
            # If position is beyond last page, return last page number
            if page_mapping and char_pos >= page_mapping[-1]["end"]:
                return page_mapping[-1]["page"]
            return None
        
        if preserve_paragraphs:
            # Split by paragraphs first
            paragraphs = text.split('\n\n')
            current_chunk = ""
            current_start = 0
            chunk_index = 0
            
            for para in paragraphs:
                para = para.strip()
                if not para:
                    continue
                
                # If adding this paragraph would exceed chunk size, save current chunk
                if current_chunk and len(current_chunk) + len(para) + 2 > chunk_size:
                    end_pos = current_start + len(current_chunk)
                    chunks.append({
                        "content": current_chunk.strip(),
                        "start_char": current_start,
                        "end_char": end_pos,
                        "chunk_index": chunk_index,
                        "page_number": find_page_number(current_start)
                    })
                    chunk_index += 1
                    
                    # Start new chunk with overlap
                    if chunk_overlap > 0 and current_chunk:
                        overlap_text = current_chunk[-chunk_overlap:]
                        current_chunk = overlap_text + "\n\n" + para
                        current_start = current_start + len(current_chunk) - len(overlap_text) - len(para) - 2
                    else:
                        current_chunk = para
                        current_start = len(text) - len(text[current_start:]) if chunks else 0
                else:
                    if current_chunk:
                        current_chunk += "\n\n" + para
                    else:
                        current_chunk = para
                        if not chunks:
                            current_start = 0
                        else:
                            # Calculate start position
                            current_start = sum(len(c["content"]) for c in chunks)
            
            # Add final chunk
            if current_chunk:
                end_pos = current_start + len(current_chunk)
                chunks.append({
                    "content": current_chunk.strip(),
                    "start_char": current_start,
                    "end_char": end_pos,
                    "chunk_index": chunk_index,
                    "page_number": find_page_number(current_start)
                })
        else:
            # Simple character-based chunking
            for i in range(0, len(text), chunk_size - chunk_overlap):
                chunk_text = text[i:i + chunk_size]
                end_pos = min(i + len(chunk_text), len(text))
                chunks.append({
                    "content": chunk_text.strip(),
                    "start_char": i,
                    "end_char": end_pos,
                    "chunk_index": len(chunks),
                    "page_number": find_page_number(i)
                })
        
        return chunks
    
    def generate_thumbnail(self, file_path: str, file_type: str, output_path: str) -> bool:
        """
        Generate a thumbnail image for a document
        
        Args:
            file_path: Path to the source document
            file_type: Type of document (pdf, jpg, png, etc.)
            output_path: Path where thumbnail should be saved
        
        Returns:
            True if thumbnail was generated successfully, False otherwise
        """
        try:
            if file_type.lower() == "pdf":
                return self._generate_pdf_thumbnail(file_path, output_path)
            elif file_type.lower() in ["jpg", "jpeg", "png", "gif", "bmp", "webp", "tiff", "tif"]:
                return self._generate_image_thumbnail(file_path, output_path)
            else:
                logger.debug(f"Thumbnail generation not supported for file type: {file_type}")
                return False
        except Exception as e:
            logger.error(f"Error generating thumbnail for {file_path}: {e}", exc_info=True)
            return False
    
    def _generate_pdf_thumbnail(self, pdf_path: str, output_path: str) -> bool:
        """Generate thumbnail from first page of PDF"""
        try:
            images = convert_from_path(pdf_path, first_page=1, last_page=1, dpi=150)
            if not images:
                return False
            
            first_page = images[0]
            thumbnail = self._resize_image(first_page, settings.THUMBNAIL_SIZE)
            thumbnail.save(output_path, "JPEG", quality=85)
            logger.info(f"Generated PDF thumbnail: {output_path}")
            return True
        except Exception as e:
            logger.error(f"Error generating PDF thumbnail: {e}", exc_info=True)
            return False
    
    def _generate_image_thumbnail(self, image_path: str, output_path: str) -> bool:
        """Generate thumbnail from image file"""
        try:
            img = Image.open(image_path)
            
            # Handle orientation (EXIF data)
            if hasattr(img, '_getexif'):
                try:
                    exif = img._getexif()
                    if exif:
                        orientation = exif.get(274)
                        if orientation == 3:
                            img = img.rotate(180, expand=True)
                        elif orientation == 6:
                            img = img.rotate(270, expand=True)
                        elif orientation == 8:
                            img = img.rotate(90, expand=True)
                except Exception:
                    pass
            
            # Convert RGBA to RGB if necessary
            if img.mode in ("RGBA", "P"):
                background = Image.new("RGB", img.size, (255, 255, 255))
                if img.mode == "P":
                    img = img.convert("RGBA")
                background.paste(img, mask=img.split()[-1] if img.mode == "RGBA" else None)
                img = background
            
            thumbnail = self._resize_image(img, settings.THUMBNAIL_SIZE)
            thumbnail.save(output_path, "JPEG", quality=85)
            logger.info(f"Generated image thumbnail: {output_path}")
            return True
        except Exception as e:
            logger.error(f"Error generating image thumbnail: {e}", exc_info=True)
            return False
    
    def _resize_image(self, img: Image.Image, max_size: int) -> Image.Image:
        """Resize image maintaining aspect ratio"""
        width, height = img.size
        if width > height:
            if width > max_size:
                new_width = max_size
                new_height = int((height * max_size) / width)
            else:
                return img
        else:
            if height > max_size:
                new_height = max_size
                new_width = int((width * max_size) / height)
            else:
                return img
        
        return img.resize((new_width, new_height), Image.Resampling.LANCZOS)
